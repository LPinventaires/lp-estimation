"""Extraction de texte des fichiers déposés + repérage des champs d'un appartement."""
import io
import re

from report import QUARTIERS


def extract_text(file_storage):
    name = (file_storage.filename or "").lower()
    data = file_storage.read()
    try:
        if name.endswith(".pdf"):
            return _pdf(data)
        if name.endswith(".docx"):
            return _docx(data)
        if name.endswith(".pptx"):
            return _pptx(data)
        return data.decode("utf-8", errors="ignore")
    except Exception:
        return ""


def _pdf(data):
    try:
        from pdfminer.high_level import extract_text as pdf_extract
        return pdf_extract(io.BytesIO(data)) or ""
    except Exception:
        try:
            from pypdf import PdfReader
            r = PdfReader(io.BytesIO(data))
            return "\n".join((p.extract_text() or "") for p in r.pages)
        except Exception:
            return ""


def _docx(data):
    from docx import Document
    d = Document(io.BytesIO(data))
    out = [p.text for p in d.paragraphs]
    for t in d.tables:
        for row in t.rows:
            out.append(" | ".join(c.text for c in row.cells))
    return "\n".join(out)


def _pptx(data):
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    out = []
    for slide in prs.slides:
        for shp in slide.shapes:
            if shp.has_text_frame:
                out.append(shp.text_frame.text)
            if shp.has_table:
                for row in shp.table.rows:
                    out.append(" | ".join(c.text for c in row.cells))
    return "\n".join(out)


def _num(s):
    s = s.replace("'", "").replace("’", "").replace(" ", "").replace(",", ".")
    try:
        return float(s)
    except ValueError:
        return None


def extract_fields(text):
    """Heuristiques de repérage. L'utilisateur valide/complète toujours ensuite."""
    f = {}
    t = text or ""
    low = t.lower()

    # Adresse : ligne "Adresse du bien : ..."
    m = re.search(r"adresse(?:\s+du\s+bien)?\s*[:\-]\s*(.+)", t, re.I)
    if m:
        f["address"] = m.group(1).strip().split("\n")[0][:160]

    # Quartier / commune
    for q in QUARTIERS:
        if q.lower() in low:
            f["quartier"] = q
            break
    m = re.search(r"commune\s*[:\-]\s*([A-Za-zÀ-ÿ' \-]+)", t, re.I)
    if m and not f.get("quartier"):
        f["quartier"] = m.group(1).strip().split("\n")[0][:80]

    # Type de bien
    for typ in ["Duplex", "Triplex", "Attique", "Appartement"]:
        if re.search(r"\b" + typ + r"\b", t, re.I):
            f["type_bien"] = typ
            break

    # Surface (m²) — "surface ... 180 m2", "X m2 PPE", "surface pondérée"
    m = (re.search(r"surface\s+pond[ée]r[ée]e[^0-9]{0,20}([0-9'’ ]{2,7})\s*m", low) or
         re.search(r"([0-9'’ ]{2,7})\s*m[²2]\s*(?:ppe|pond)", low) or
         re.search(r"d['’]une\s+surface[^0-9]{0,12}([0-9'’ ]{2,7})\s*m", low) or
         re.search(r"surface[^0-9]{0,12}([0-9'’ ]{2,7})\s*m[²2]", low))
    if m:
        v = _num(m.group(1))
        if v and 15 < v < 2000:
            f["surface"] = v

    # Pièces
    m = re.search(r"([0-9]+(?:[\.,][0-9])?)\s*pi[èe]ces", low)
    if m:
        f["pieces"] = m.group(1).replace(",", ".")

    # Étage
    m = re.search(r"(\d+)\s*(?:er|ème|e)\s*[ée]tage", low)
    if m:
        f["etage"] = m.group(1)
    elif "rez" in low:
        f["etage"] = "Rez"

    # Année de construction
    m = re.search(r"(?:construction|construit)[^0-9]{0,20}((?:~\s*)?(?:19|20)\d{2})", low)
    if m:
        f["annee"] = m.group(1).replace(" ", "")
    elif re.search(r"avant\s+1919|ant[ée]rieur\s+[àa]\s+1919", low):
        f["annee"] = "avant 1919"

    # Parkings
    m = re.search(r"(\d+)\s*(?:places?\s+de\s+park|place\s+de\s+parc|box)", low)
    if m:
        f["parking_nb"] = int(m.group(1))
    elif "deux places de parking" in low or "2 places" in low:
        f["parking_nb"] = 2

    # Balcon / loggia / terrasse
    m = re.search(r"(?:balcon|loggia|terrasse)[^0-9]{0,18}([0-9'’ ]{1,5})\s*m", low)
    if m:
        v = _num(m.group(1))
        if v and v < 300:
            f["balcon"] = v

    # État
    if "excellent" in low or "parfait état" in low:
        f["etat"] = "Excellent"
    elif "bon état" in low:
        f["etat"] = "Bon"
    elif "rénové" in low or "renove" in low:
        f["etat"] = "Rénové"
    elif "à rafraîchir" in low or "travaux" in low:
        f["etat"] = "Travaux à prévoir"

    return f
